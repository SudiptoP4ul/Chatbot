import os
import re
import sqlite3
import glob
import numpy as np
import logging
from difflib import get_close_matches
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

logging.getLogger('zeep').setLevel(logging.ERROR)

BASE_DIR     = os.path.dirname(os.path.abspath(__file__))
DB_PATH      = os.path.join(BASE_DIR, 'swr_contingency.db')
STATIONS_CSV = os.path.join(BASE_DIR, 'stations.csv')

PPTM_FILES = [
    os.path.join(BASE_DIR, 'Mainline_CPT__8th_May_21_.pptm'),
    os.path.join(BASE_DIR, '20_Batch_2_LOR_2-4_CPT__19th_May_21_.pptm'),
]

DOCX_DIRS = [
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM West'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM Central'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM Metro'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM South'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM Waterloo'),
]

KNOWN_LABELS = [
    'Principles of Service Alteration',
    'Alternative Passenger Journey',
    'Information for Signallers',
    'Additional Information for Station Staff',
    'Passenger Information',
]

SECTION_KEYWORDS = ['off peak', 'peak', 'during peak', 'all day', 'note', 'if no access']

SWR_CONTACTS = [
    ('Bournemouth Main Line', 'Basingstoke SCP',               '50297'),
    ('Bournemouth Main Line', 'Basingstoke ASC',               '0750131'),
    ('Bournemouth Main Line', 'Basingstoke (DC) Panel',        '0750010'),
    ('Bournemouth Main Line', 'Woking SCP',                    '50341'),
    ('Bournemouth Main Line', 'Woking ASC',                    '0069410 / 0069420'),
    ('Bournemouth Main Line', 'Southampton SCP',               '075 8219'),
    ('Bournemouth Main Line', 'Southampton Ctl Supv',          '075 8107'),
    ('Bournemouth Main Line', 'Eastleigh ASC',                 '0757202'),
    ('Bournemouth Main Line', 'Bournemouth SCP',               '50375'),
    ('Bournemouth Main Line', 'Bournemouth SB',                '0755332'),
    ('Bournemouth Main Line', 'Brockenhurst SB',               '08528311'),
    ('Bournemouth Main Line', 'Guildford SB',                  '0067562'),
    ('Bournemouth Main Line', 'Guildford SCP',                 '50359'),
    ('Bournemouth Main Line', 'Portsmouth SCP',                '50265'),
    ('Bournemouth Main Line', 'Dorchester SB',                 '0758355'),
    ('Bournemouth Main Line', 'Salisbury SCP',                 '50412'),
    ('Bournemouth Main Line', 'GWR Train Service Control West','08582219'),
    ('Bournemouth Main Line', 'Guildford Surbiton Supervisor', '07771 613061'),
]

TERMINOLOGY = {
    'DM':   'Down Main - direction of travel towards London',
    'UM':   'Up Main - direction of travel away from London',
    'SL':   'Slow Line',
    'FL':   'Fast Line',
    'DMSL': 'Down Main Slow Line - slow line towards London',
    'DMFL': 'Down Main Fast Line - fast line towards London',
    'UMSL': 'Up Main Slow Line - slow line away from London',
    'UMFL': 'Up Main Fast Line - fast line away from London',
    'SWR':  'South Western Railway',
    'NR':   'Network Rail',
    'TOC':  'Train Operating Company',
    'ASC':  'Area Signalling Centre',
    'SCP':  'Station Control Point',
    'SB':   'Signal Box',
    'CSL2': 'Customer Service Level 2 - disruption response protocol',
    'CRS':  'Computer Reservation System - 3-letter station code',
}

# intent detection keyword sets
BLOCKAGE_WORDS    = ['blockage','blocked','block','line blocked','track blocked',
                     'dmfl','dmsl','umfl','umsl','down fast','down slow',
                     'up fast','up slow','full blockage','partial blockage',
                     'one line','both lines','derailment','obstruction']

DISRUPTION_WORDS  = ['station','disruption','disrupted','closed','closure',
                     'station problem','not stopping','trains not stopping',
                     'station disruption','station closed','platform']

PASSENGER_WORDS   = ['passenger','passengers','alternative','route','routes',
                     'bus','coach','taxi','travel','get to','how to get',
                     'alternative route','what can passengers','advice for passenger']

CONTACT_WORDS     = ['contact','phone','number','call','who do i call',
                     'telephone','ring','escalate','control']

STAFF_WORDS       = ['staff','signaller','signal','what should i do',
                     'what do i do','instructions','principles','alteration',
                     'service alteration','divert','diversion']

TERM_WORDS        = ['what is','what does','what does it mean','explain',
                     'meaning','mean','define']


def get_blockage_type_id(status):
    sl = status.lower()
    if 'dmfl' in sl and 'dmsl' in sl: return 5
    if 'umfl' in sl and 'umsl' in sl: return 5
    if 'dmfl' in sl: return 1
    if 'dmsl' in sl: return 2
    if 'umfl' in sl: return 3
    if 'umsl' in sl: return 4
    if 'both' in sl or ('all' in sl and 'line' in sl): return 5
    if 'one line' in sl or 'partial' in sl: return 6
    if 'station' in sl: return 7
    return 6


def detect_section(line):
    ll = line.lower().strip()
    for kw in SECTION_KEYWORDS:
        if kw in ll:
            return line.strip()
    return ''


def is_important(line):
    ll = line.lower()
    return any(x in ll for x in
               ['must','immediately','do not','cancel','emergency','warning'])


def build_database(db_path=DB_PATH):
    try:
        from pptx import Presentation
        from docx import Document
    except ImportError:
        os.system("pip install python-pptx python-docx -q")
        from pptx import Presentation
        from docx import Document

    if os.path.exists(db_path):
        os.remove(db_path)

    conn = sqlite3.connect(db_path)
    cur  = conn.cursor()

    cur.executescript("""
    CREATE TABLE Tiploc (
        name TEXT, longname TEXT, name_alias TEXT,
        alpha3 TEXT, tiploc TEXT PRIMARY KEY
    );
    CREATE TABLE blockage_types (id INTEGER PRIMARY KEY, description TEXT);
    CREATE TABLE ContingencyPlan (
        id INTEGER, subplan INTEGER,
        name_of_line TEXT, first_station TEXT, last_station TEXT,
        blockage_type INTEGER, blockage_status TEXT,
        PRIMARY KEY (id, subplan)
    );
    CREATE TABLE ServiceAlteration (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, alteration TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE AltPassengerJourney (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, altjourney TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE SignallerInfo (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, info TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE StationStaffInfo (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, info TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE PassengerInfo (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, info TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE ContactDetails (
        name_of_line TEXT, contact_name TEXT, phone_numbers TEXT
    );
    CREATE TABLE StationDisruptionPlan (
        id INTEGER PRIMARY KEY, station_name TEXT,
        dateof TEXT, doc_owner TEXT, region TEXT
    );
    CREATE TABLE OtherStationSupport (
        dplan_id INTEGER, other_station_name TEXT,
        phone_numbers TEXT, time_info TEXT
    );
    CREATE TABLE AltTransportInfo (
        dplan_id INTEGER, destination TEXT,
        route_type INTEGER, route_info TEXT
    );
    CREATE TABLE StationTips (dplan_id INTEGER, issue TEXT, tip TEXT);
    """)

    cur.executemany("INSERT INTO blockage_types VALUES (?,?)", [
        (1,"DMFL - Down Main Fast Line blocked"),
        (2,"DMSL - Down Main Slow Line blocked"),
        (3,"UMFL - Up Main Fast Line blocked"),
        (4,"UMSL - Up Main Slow Line blocked"),
        (5,"Both lines blocked (full blockage)"),
        (6,"One line blocked (partial blockage)"),
        (7,"Station blocked or closed"),
    ])

    try:
        import pandas as pd
        if os.path.exists(STATIONS_CSV):
            df = pd.read_csv(STATIONS_CSV).reset_index()
            df.columns = ['name','longname','name_alias','alpha3','tiploc']
            df['name_alias'] = df['name_alias'].replace('\\N','')
            for _, row in df.iterrows():
                cur.execute("INSERT OR IGNORE INTO Tiploc VALUES (?,?,?,?,?)",
                    (row['name'],row['longname'],row['name_alias'],
                     row['alpha3'],row['tiploc']))
            print(f"Loaded {len(df)} stations into Tiploc table")
        else:
            print("Warning: stations.csv not found")
    except Exception as e:
        print(f"Warning: could not load stations.csv: {e}")

    cur.executemany("INSERT INTO ContactDetails VALUES (?,?,?)", SWR_CONTACTS)

    plan_count = 0
    for filepath in PPTM_FILES:
        if not os.path.exists(filepath):
            print(f"Warning: {filepath} not found, skipping.")
            continue
        prs = Presentation(filepath)
        for slide in prs.slides:
            title = ""
            table_cells = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape,"text") and shape.text.strip():
                        if not title: title = shape.text.strip()
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                t = cell.text.strip()
                                if t: table_cells.append(t)
                except Exception:
                    pass

            if not re.match(r'Plan\s+[\d\.a]+', title):
                continue

            title_clean = title.replace('\t',' ').replace('\x0b',' ').strip()
            m = re.match(r'Plan\s+([\d]+)[\.a]?([\d]*)\s+(.*)', title_clean)
            if not m: continue

            plan_id     = int(m.group(1))
            subplan     = int(m.group(2)) if m.group(2) else 0
            loc_str     = m.group(3).strip()
            parts       = re.split(r'\s+-\s+', loc_str)
            first_st    = parts[0].strip() if parts else ""
            last_st     = parts[1].strip() if len(parts) > 1 else ""

            line = ""; status = ""; field_map = {}
            for idx, cell in enumerate(table_cells):
                if ('Main Line' in cell or 'Line' in cell) and len(cell) < 60:
                    line = cell.replace('\n',' ').strip()
                elif cell.startswith('Status:'):
                    status = cell.replace('Status:','').strip()
                elif cell in KNOWN_LABELS and idx+1 < len(table_cells):
                    field_map[cell] = table_cells[idx+1]

            cur.execute(
                "INSERT OR IGNORE INTO ContingencyPlan VALUES (?,?,?,?,?,?,?)",
                (plan_id,subplan,line,first_st,last_st,
                 get_blockage_type_id(status),status))

            def insert_bullets(table, col, text, pid, sp):
                if not text or text == 'N/A': return
                current_section = ''
                for ln in text.split('\n'):
                    ln = ln.strip().lstrip('•').strip()
                    if not ln: continue
                    sec = detect_section(ln)
                    if sec:
                        current_section = sec
                        continue
                    imp = 1 if is_important(ln) else 0
                    cur.execute(f"INSERT INTO {table} VALUES (?,?,?,?,?)",
                                (pid,sp,current_section,ln,imp))

            insert_bullets("ServiceAlteration","alteration",
                field_map.get('Principles of Service Alteration',''),plan_id,subplan)
            insert_bullets("AltPassengerJourney","altjourney",
                field_map.get('Alternative Passenger Journey',''),plan_id,subplan)
            insert_bullets("SignallerInfo","info",
                field_map.get('Information for Signallers',''),plan_id,subplan)
            insert_bullets("StationStaffInfo","info",
                field_map.get('Additional Information for Station Staff',''),plan_id,subplan)
            insert_bullets("PassengerInfo","info",
                field_map.get('Passenger Information',''),plan_id,subplan)
            plan_count += 1

    print(f"Loaded {plan_count} contingency plans from PowerPoint files")

    docx_files = []
    for d in DOCX_DIRS:
        if os.path.exists(d):
            docx_files += glob.glob(os.path.join(d,"*.docx"))
    docx_files += glob.glob(os.path.join(BASE_DIR,"Station_Disruption_Plan_*.docx"))
    docx_files  = sorted(set(docx_files))

    dplan_id = 1; station_count = 0
    for filepath in docx_files:
        try:
            doc = Document(filepath)
            station_name = ""; doc_owner = ""; dateof = ""
            region = "Unknown"
            for r in ["RM West","RM Central","RM Metro","RM South","RM Waterloo"]:
                if r in filepath: region = r; break
            if region == "Unknown": region = "RM West"

            for para in doc.paragraphs:
                text = para.text.strip()
                if text and para.style.name.startswith("Heading") and not station_name:
                    station_name = text.replace("Station","").strip()
                    break

            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if len(cells) >= 2:
                        if "Document Owner" in cells[0]: doc_owner = cells[1]
                        if "Date" in cells[0] and not dateof: dateof = cells[1]

            if not station_name:
                fn = os.path.basename(filepath)
                m2 = re.search(r'Plan_-_(.+?)_Issue', fn)
                if m2: station_name = m2.group(1).replace('_',' ').strip()

            cur.execute("INSERT INTO StationDisruptionPlan VALUES (?,?,?,?,?)",
                        (dplan_id,station_name,dateof,doc_owner,region))

            for table in doc.tables:
                headers = [c.text.strip().upper() for c in table.rows[0].cells]
                if any("ROUTE AVAILABLE" in h for h in headers):
                    for row in table.rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        dest  = cells[1] if len(cells) > 1 else ""
                        if not dest or dest in ("Destination",""): continue
                        rail    = cells[0] if len(cells) > 0 else ""
                        bus     = cells[2] if len(cells) > 2 else ""
                        walking = cells[3] if len(cells) > 3 else ""
                        if rail: cur.execute("INSERT INTO AltTransportInfo VALUES (?,?,?,?)",
                                             (dplan_id,dest,0,rail))
                        if bus:  cur.execute("INSERT INTO AltTransportInfo VALUES (?,?,?,?)",
                                             (dplan_id,dest,1,bus))
                        if walking and walking.lower() not in ('n/a',''):
                            cur.execute("INSERT INTO AltTransportInfo VALUES (?,?,?,?)",
                                        (dplan_id,dest,2,walking))

                if len(headers) >= 2 and "ISSUE" in headers[0]:
                    for row in table.rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        if cells[0] and cells[0].upper() not in ("ISSUE",""):
                            cur.execute("INSERT INTO StationTips VALUES (?,?,?)",
                                (dplan_id,cells[0][:255],
                                 cells[1][:500] if len(cells)>1 else ""))

                if any("STATION NAME" in h or "TELEPHONE" in h for h in headers):
                    for row in table.rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        if cells[0] and cells[0] not in ("N/A","Station Name",""):
                            cur.execute("INSERT INTO OtherStationSupport VALUES (?,?,?,?)",
                                (dplan_id,cells[0],
                                 cells[1] if len(cells)>1 else "",
                                 cells[2] if len(cells)>2 else ""))

            dplan_id += 1; station_count += 1
        except Exception as e:
            print(f"Error reading {os.path.basename(filepath)}: {e}")

    print(f"Loaded {station_count} station disruption plans")
    conn.commit(); conn.close()
    print("Database build complete.")


class ContingencyDB:
    def __init__(self, db_path=DB_PATH):
        if not os.path.exists(db_path):
            print("Database not found. Building now...")
            build_database(db_path)
        self.conn = sqlite3.connect(db_path)
        self.conn.row_factory = sqlite3.Row
        self._build_search_index()

    def _build_search_index(self):
        cur = self.conn.cursor()
        cur.execute("SELECT id,subplan,first_station,last_station,blockage_status FROM ContingencyPlan")
        rows = cur.fetchall()
        self.cp_ids    = [(r['id'],r['subplan']) for r in rows]
        self.cp_corpus = [
            f"{r['first_station']} {r['last_station']} {r['blockage_status']}"
            for r in rows
        ]
        for i,(pid,sp) in enumerate(self.cp_ids):
            cur.execute("SELECT alteration FROM ServiceAlteration WHERE cplan_id=? AND subplan=?",(pid,sp))
            self.cp_corpus[i] += " " + " ".join(r['alteration'] for r in cur.fetchall())

        self.cp_vectorizer = TfidfVectorizer(stop_words='english',ngram_range=(1,2),sublinear_tf=True)
        self.cp_matrix     = self.cp_vectorizer.fit_transform(self.cp_corpus)

        cur.execute("SELECT id,station_name FROM StationDisruptionPlan")
        rows          = cur.fetchall()
        self.dp_ids   = [r['id']           for r in rows]
        self.dp_names = [r['station_name'] for r in rows]

        self.cp_stations = sorted(set(
            s[0] for s in self.conn.execute("SELECT first_station FROM ContingencyPlan") if s[0]
        ) | set(
            s[0] for s in self.conn.execute("SELECT last_station FROM ContingencyPlan") if s[0]
        ))
        self.dp_stations = sorted(set(s for s in self.dp_names if s))

        self.tiploc_cache = {}
        for row in self.conn.execute("SELECT name,alpha3,tiploc FROM Tiploc"):
            self.tiploc_cache[row['name'].upper()] = {
                'crs': row['alpha3'], 'tiploc': row['tiploc']}

    def fuzzy_match(self, raw, station_list):
        if not raw: return None
        raw = raw.strip().title()
        if raw in station_list: return raw
        close = get_close_matches(raw, station_list, n=1, cutoff=0.5)
        if close: return close[0]
        for s in station_list:
            if raw.lower() in s.lower(): return s
        return None

    def extract_station_from_text(self, text, station_list):
        text_title = text.strip().title()
        for station in station_list:
            if station.lower() in text.lower():
                return station
        return self.fuzzy_match(text_title, station_list)

    def search_contingency(self, query, loc_from='', loc_to='', top_n=3):
        full_query = f"{loc_from} {loc_to} {query}".strip()
        q_vec      = self.cp_vectorizer.transform([full_query])
        scores     = cosine_similarity(q_vec, self.cp_matrix)[0]
        loc_from_m = self.fuzzy_match(loc_from, self.cp_stations) if loc_from else None
        loc_to_m   = self.fuzzy_match(loc_to,   self.cp_stations) if loc_to   else None
        cur        = self.conn.cursor()
        boosted    = []
        for i,(pid,sp) in enumerate(self.cp_ids):
            score = float(scores[i])
            cur.execute("SELECT first_station,last_station FROM ContingencyPlan WHERE id=? AND subplan=?",(pid,sp))
            row = cur.fetchone()
            if row:
                if loc_from_m and row['first_station'].lower() == loc_from_m.lower(): score += 0.4
                if loc_to_m   and row['last_station'].lower()  == loc_to_m.lower():   score += 0.4
            boosted.append((pid,sp,score))
        boosted.sort(key=lambda x: x[2], reverse=True)
        results = []
        for pid,sp,score in boosted[:top_n]:
            if score > 0.01:
                plan = self.get_full_plan(pid,sp,score)
                if plan: results.append(plan)
        return results

    def search_station_disruption(self, station_name):
        matched = self.fuzzy_match(station_name, self.dp_stations)
        if matched:
            cur = self.conn.cursor()
            cur.execute("SELECT id FROM StationDisruptionPlan WHERE station_name=?",(matched,))
            row = cur.fetchone()
            if row: return self.get_disruption_plan(row['id'])
        return None

    def get_full_plan(self, plan_id, subplan, score=0):
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM ContingencyPlan WHERE id=? AND subplan=?",(plan_id,subplan))
        cp = cur.fetchone()
        if not cp: return None
        def get_bullets(table, col):
            cur.execute(
                f"SELECT section,{col},important FROM {table} "
                f"WHERE cplan_id=? AND subplan=? ORDER BY rowid",(plan_id,subplan))
            return [{'text':r[1],'section':r[0],'important':bool(r[2])} for r in cur.fetchall()]
        return {
            "plan_id":        plan_id,  "subplan":        subplan,
            "score":          score,    "line":           cp['name_of_line'],
            "first_station":  cp['first_station'],
            "last_station":   cp['last_station'],
            "blockage_status": cp['blockage_status'],
            "principles":     get_bullets("ServiceAlteration",  "alteration"),
            "alt_journey":    get_bullets("AltPassengerJourney","altjourney"),
            "signaller":      get_bullets("SignallerInfo",       "info"),
            "staff":          get_bullets("StationStaffInfo",   "info"),
            "passenger":      get_bullets("PassengerInfo",      "info"),
        }

    def get_disruption_plan(self, dplan_id):
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM StationDisruptionPlan WHERE id=?",(dplan_id,))
        dp = cur.fetchone()
        if not dp: return None
        cur.execute(
            "SELECT destination,route_type,route_info FROM AltTransportInfo "
            "WHERE dplan_id=? ORDER BY destination,route_type",(dplan_id,))
        transport = {}
        for row in cur.fetchall():
            dest = row['destination']
            if dest not in transport:
                transport[dest] = {'destination':dest,'rail':'','bus':'','walking':''}
            if   row['route_type'] == 0: transport[dest]['rail']    = row['route_info']
            elif row['route_type'] == 1: transport[dest]['bus']     = row['route_info']
            elif row['route_type'] == 2: transport[dest]['walking'] = row['route_info']
        cur.execute("SELECT * FROM StationTips           WHERE dplan_id=?",(dplan_id,))
        tips    = [dict(r) for r in cur.fetchall()]
        cur.execute("SELECT * FROM OtherStationSupport  WHERE dplan_id=?",(dplan_id,))
        support = [dict(r) for r in cur.fetchall()]
        return {
            "station_name": dp['station_name'], "dateof":    dp['dateof'],
            "region":       dp['region'],       "doc_owner": dp['doc_owner'],
            "transport":    list(transport.values()),
            "tips": tips, "support": support,
        }

    def get_contacts(self, line_name=None):
        cur = self.conn.cursor()
        if line_name:
            cur.execute("SELECT contact_name,phone_numbers FROM ContactDetails WHERE name_of_line LIKE ?",(f'%{line_name}%',))
        else:
            cur.execute("SELECT contact_name,phone_numbers FROM ContactDetails")
        return [dict(r) for r in cur.fetchall()]

    def lookup_tiploc(self, station_name):
        return self.tiploc_cache.get(station_name.upper().strip(), None)

    def get_all_cp_stations(self): return self.cp_stations
    def get_all_dp_stations(self): return self.dp_stations
    def explain_term(self, term):  return TERMINOLOGY.get(term.upper().strip(), None)


def contains(text, words):
    t = text.lower()
    return any(w in t for w in words)


def extract_station(text, db):
    all_stations = db.get_all_cp_stations() + db.get_all_dp_stations()
    for station in all_stations:
        if station.lower() in text.lower():
            return station
    words = text.replace(',','').split()
    for i in range(len(words)):
        for j in range(i+1, min(i+4, len(words)+1)):
            candidate = ' '.join(words[i:j])
            matched   = db.fuzzy_match(candidate, all_stations)
            if matched:
                return matched
    return None


def bot(msg):
    print(f"Bot: {msg}")


def ask(prompt):
    return input(f"Bot: {prompt}\nYou: ").strip()


def print_bullets(bullets):
    current_section = ''
    for b in bullets:
        if b['section'] and b['section'] != current_section:
            current_section = b['section']
            print(f"     [{current_section}]")
        prefix = "  ! " if b['important'] else "  - "
        print(f"{prefix}{b['text']}")


def show_passenger_advice(plan, source='contingency'):
    if source == 'disruption':
        if plan['transport']:
            bot("Here are the alternative travel options for passengers:")
            print()
            for t in plan['transport']:
                if t['destination']:
                    print(f"  To {t['destination']}:")
                    if t['rail']:    print(f"    Rail    - {t['rail']}")
                    if t['bus']:     print(f"    Bus     - {t['bus']}")
                    if t['walking']: print(f"    Walking - {t['walking']}")
            print()
        if plan['tips']:
            bot("Some useful tips for this station:")
            for t in plan['tips']:
                if t['issue']:
                    print(f"  - {t['issue'][:100]}")
    else:
        if plan['alt_journey']:
            bot("Here is the alternative passenger journey advice:")
            print()
            print_bullets(plan['alt_journey'])
            print()
        if plan['passenger']:
            bot("Detailed passenger information:")
            print()
            print_bullets(plan['passenger'])
            print()


def show_staff_advice(plan):
    if plan['principles']:
        bot("Principles of service alteration for this incident:")
        print()
        print_bullets(plan['principles'])
        print()
    if plan['staff']:
        bot("Instructions for station staff:")
        print()
        print_bullets(plan['staff'])
        print()


def show_signaller_advice(plan):
    if plan['signaller']:
        bot("Information for signallers:")
        print()
        print_bullets(plan['signaller'])
        print()


def show_contacts(db):
    contacts = db.get_contacts()
    bot("Here are the control contact numbers:")
    print()
    for c in contacts:
        print(f"  {c['contact_name']:<40} {c['phone_numbers']}")
    print()


def show_full_plan(plan, role='all'):
    print()
    print(f"  Plan {plan['plan_id']}.{plan['subplan']} | "
          f"{plan['first_station']} to {plan['last_station']}")
    print(f"  Status: {plan['blockage_status']}")
    print(f"  Line  : {plan['line']}")
    print()
    if role in ('all','staff'):
        show_staff_advice(plan)
    if role in ('all','signaller'):
        show_signaller_advice(plan)
    if role in ('all','passenger'):
        show_passenger_advice(plan)


#  MAIN CHATBOT CONVERSATION

_db = None

def get_db():
    global _db
    if _db is None:
        _db = ContingencyDB(DB_PATH)
    return _db


def contingency_flow():
    db = get_db()

    print()
    bot("Hi, I'm the Contingency Assistant. I'm here to help you "
        "manage disruptions and find the right information quickly.")
    bot("Tell me what's happening and I'll find the right plan for you. "
        "You can describe the situation in your own words.")
    print()

    user_input = input("You: ").strip()
    if not user_input:
        bot("I didn't catch that. Please describe the incident.")
        return

    _handle_input(db, user_input)


def _handle_input(db, user_input):
    text = user_input.lower()

    # check for terminology question first
    for term, explanation in TERMINOLOGY.items():
        if term.lower() in text and contains(text, TERM_WORDS + ['what is', 'what does']):
            bot(f"{term} stands for: {explanation}")
            _ask_anything_else(db)
            return

    # check for contact request
    if contains(text, CONTACT_WORDS):
        show_contacts(db)
        _ask_anything_else(db)
        return

    # try to extract a station name from what they said
    station = extract_station(user_input, db)

    # decide incident type from keywords
    is_blockage   = contains(text, BLOCKAGE_WORDS)
    is_disruption = contains(text, DISRUPTION_WORDS)

    # if both or neither, ask to clarify
    if is_blockage and not is_disruption:
        _handle_blockage(db, user_input, station)
        return

    if is_disruption and not is_blockage:
        _handle_disruption(db, user_input, station)
        return

    # ambiguous - ask what type
    bot("I can see there's an incident. Is this:")
    print("  1) A line blockage (track blocked between two stations)")
    print("  2) A station disruption (problems at a specific station)")
    choice = input("\nYou: ").strip()

    if choice == '2' or contains(choice, DISRUPTION_WORDS + ['station','2']):
        _handle_disruption(db, user_input, station)
    else:
        _handle_blockage(db, user_input, station)


def _handle_blockage(db, user_input, station_hint=None):
    bot("Ok, I'll help you find the contingency plan for this line blockage.")

    # get location from
    if station_hint and station_hint in db.get_all_cp_stations():
        loc_from = station_hint
        bot(f"I can see this is near {loc_from}. Is that the start of the blockage?")
        confirm = input("You: ").strip().lower()
        if confirm not in ('','y','yes'):
            loc_from = _ask_station(db, "Which station is the START of the blockage?",
                                    db.get_all_cp_stations())
    else:
        loc_from = _ask_station(db, "Which station is the START of the blockage?",
                                db.get_all_cp_stations())

    if not loc_from:
        bot("I couldn't identify the start station. Please contact Network Control.")
        show_contacts(db)
        return

    # get location to
    loc_to = ""
    ans = ask("Which station is the END of the blockage? (or press Enter if you're not sure)")
    if ans.strip():
        loc_to = db.fuzzy_match(ans.strip(), db.get_all_cp_stations()) or ""
        if loc_to and loc_to.lower() != ans.lower():
            bot(f"Got it, I'll look for plans up to {loc_to}.")

    # get blockage type
    blockage_desc = ask("What type of blockage is it? "
                        "(e.g. DMFL, down slow line, full blockage, one line blocked)")
    if not blockage_desc:
        blockage_desc = "blockage"

    # explain term if they typed a code
    first_word = blockage_desc.split()[0].upper()
    term = db.explain_term(first_word)
    if term:
        bot(f"({first_word} = {term})")

    # search
    bot("Searching the knowledge base...")
    results = db.search_contingency(blockage_desc, loc_from, loc_to, top_n=3)

    if not results:
        bot(f"I couldn't find a specific contingency plan for this incident.")
        bot("Please escalate to Network Control immediately.")
        show_contacts(db)
        return

    best = results[0]
    confidence = "High" if best['score'] > 0.5 else "Medium" if best['score'] > 0.2 else "Low"
    bot(f"I found {len(results)} matching plan(s). Best match is Plan "
        f"{best['plan_id']}.{best['subplan']} ({confidence} confidence).")
    print(f"  Location : {best['first_station']} to {best['last_station']}")
    print(f"  Status   : {best['blockage_status']}")
    print(f"  Line     : {best['line']}")
    print()

    # ask what they need
    bot("What information do you need from this plan?")
    print("  1) Passenger advice and alternative routes")
    print("  2) Instructions for station staff")
    print("  3) Information for signallers")
    print("  4) Everything")
    info_choice = input("\nYou: ").strip().lower()

    if info_choice == '1' or contains(info_choice, PASSENGER_WORDS):
        show_passenger_advice(best, 'contingency')
    elif info_choice == '2' or contains(info_choice, STAFF_WORDS):
        show_staff_advice(best)
    elif info_choice == '3' or 'signal' in info_choice:
        show_signaller_advice(best)
    else:
        show_full_plan(best, 'all')

    # offer more plans
    if len(results) > 1:
        more = ask(f"There are {len(results)-1} more related plan(s). "
                   f"Would you like to see them? (y/n)")
        if more.lower() in ('y','yes'):
            for plan in results[1:]:
                show_full_plan(plan, 'all')

    # show tiploc info in summary
    tiploc_from = db.lookup_tiploc(loc_from.upper())
    tiploc_to   = db.lookup_tiploc(loc_to.upper()) if loc_to else None

    print()
    bot("Incident summary:")
    from_line = f"  Start : {loc_from}"
    if tiploc_from: from_line += f" (CRS: {tiploc_from['crs']})"
    print(from_line)
    if loc_to:
        to_line = f"  End   : {loc_to}"
        if tiploc_to: to_line += f" (CRS: {tiploc_to['crs']})"
        print(to_line)
    print(f"  Type  : {blockage_desc}")
    print(f"  Plan  : Plan {best['plan_id']}.{best['subplan']}")
    print()

    _ask_anything_else(db)


def _handle_disruption(db, user_input, station_hint=None):
    bot("Ok, I'll look up the station disruption plan.")

    # get station
    if station_hint and station_hint in db.get_all_dp_stations():
        station = station_hint
        bot(f"Is the disruption at {station}?")
        confirm = input("You: ").strip().lower()
        if confirm not in ('', 'y', 'yes'):
            station = _ask_station(db, "Which station is affected?",
                                   db.get_all_dp_stations())
    else:
        station = _ask_station(db, "Which station is affected?",
                               db.get_all_dp_stations())

    if not station:
        bot("I couldn't identify the station. Please contact Network Control.")
        show_contacts(db)
        return

    bot(f"Looking up the disruption plan for {station}...")
    plan = db.search_station_disruption(station)

    if not plan:
        bot(f"I don't have a specific disruption plan for {station}.")
        bot("Please contact Network Control for assistance.")
        show_contacts(db)
        return

    bot(f"I found the disruption plan for {station}. What do you need?")
    print("  1) Alternative routes for passengers")
    print("  2) Station tips and advice")
    print("  3) Nearby station contacts")
    print("  4) Everything")
    info_choice = input("\nYou: ").strip().lower()

    if info_choice == '1' or contains(info_choice, PASSENGER_WORDS):
        show_passenger_advice(plan, 'disruption')

    elif info_choice == '2' or 'tip' in info_choice or 'advice' in info_choice:
        if plan['tips']:
            bot(f"Here are the top tips for {station}:")
            print()
            for t in plan['tips']:
                if t['issue']:
                    print(f"  Issue : {t['issue'][:100]}")
                    if t['tip']:
                        print(f"  Tip   : {t['tip'][:100]}")
            print()
        else:
            bot("No specific tips found for this station.")

    elif info_choice == '3' or contains(info_choice, CONTACT_WORDS):
        if plan['support']:
            bot("Nearby station support contacts:")
            print()
            for s in plan['support']:
                if s['other_station_name']:
                    line = f"  {s['other_station_name']}"
                    if s['phone_numbers']: line += f" - Tel: {s['phone_numbers']}"
                    print(line)
            print()
        else:
            bot("No nearby station contacts found in this plan.")
            bot("Here are the main control numbers instead:")
            show_contacts(db)

    else:
        # show everything
        bot(f"Full disruption plan for {station}:")
        print(f"  Region  : {plan['region']}")
        if plan['dateof']:    print(f"  Date    : {plan['dateof']}")
        if plan['doc_owner']: print(f"  Owner   : {plan['doc_owner']}")
        print()
        show_passenger_advice(plan, 'disruption')
        if plan['tips']:
            bot("Station tips:")
            for t in plan['tips']:
                if t['issue']:
                    print(f"  - {t['issue'][:100]}")
            print()
        if plan['support']:
            bot("Nearby station contacts:")
            for s in plan['support']:
                if s['other_station_name']:
                    print(f"  {s['other_station_name']} - Tel: {s['phone_numbers']}")
            print()

    tiploc = db.lookup_tiploc(station.upper())
    if tiploc:
        bot(f"Station CRS code: {tiploc['crs']} | TIPLOC: {tiploc['tiploc']}")

    _ask_anything_else(db)


def _ask_station(db, prompt, station_list):
    while True:
        raw = input(f"\nBot: {prompt}\nYou: ").strip()
        if not raw:
            bot("Please enter a station name.")
            continue
        if raw.lower() == 'list':
            bot(f"Known stations ({len(station_list)} total):")
            for s in station_list:
                print(f"  {s}")
            continue
        term = db.explain_term(raw)
        if term:
            bot(f"{raw.upper()} means: {term}")
            continue
        matched = db.fuzzy_match(raw, station_list)
        if matched:
            if matched.lower() != raw.lower():
                ok = input(f"Bot: Did you mean {matched}? [Y/n]: ").strip().lower()
                if ok not in ('','y','yes'):
                    bot("Please try again. Type list to see all known stations.")
                    continue
            return matched
        bot(f"I couldn't find '{raw}'. Type list to see all known stations.")
        again = input("Bot: Try again? [Y/n]: ").strip().lower()
        if again not in ('','y','yes'):
            return None


def _ask_anything_else(db):
    print()
    ans = ask("Is there anything else I can help you with? (y/n)")
    if ans.lower() in ('y','yes'):
        bot("Ok, what else do you need?")
        user_input = input("You: ").strip()
        _handle_input(db, user_input)
    else:
        bot("Ok, good luck managing the incident. Stay safe.")


def main():
    print("Initializing Contingency Expert System...")
    get_db()

    while True:
        print("\n STAFF MENU")
        print("1) Contingency Assistant")
        print("2) Exit")

        choice = input("> ").strip()
        if choice == '1':
            contingency_flow()
        elif choice == '2':
            break
        else:
            print("Please select 1 or 2.")


if __name__ == "__main__":
    main()
