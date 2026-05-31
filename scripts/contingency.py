import os
import glob
import sqlite3
import datetime
import re
import pandas as pd

try:
    from pptx import Presentation
    from docx import Document
except ImportError:
    Presentation = None
    Document = None

BASE_DIR = '/Users/sudiptogoldfish/code files/7059B A_AI Lab/Chatbot/'
KB_DIR = os.path.join(BASE_DIR, 'knowledge based') 
DATA_DIR = os.path.join(BASE_DIR, 'train service data')
DB_PATH = os.path.join(DATA_DIR, 'integrated_system.db')

class ContingencyExpertSystem:
    def __init__(self):
        os.makedirs(DATA_DIR, exist_ok=True)
        self.init_sqlite_tables()
        self.stations_df = self.sync_and_load_stations()
        self.mined_knowledge_pool = []
        self.load_and_mine_local_docs()

    def init_sqlite_tables(self):
        """Initializes system tables and safely executes cleaned user SQL schema layout loops."""
        conn = sqlite3.connect(DB_PATH)
        cursor = conn.cursor()
        
        # 1. Framework Tracking System Core
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS conversation_history (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp TEXT,
                session_mode TEXT,
                user_message TEXT,
                bot_response TEXT
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS incident_logs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp TEXT,
                location TEXT,
                event_type TEXT,
                severity TEXT,
                staff_advice TEXT,
                passenger_advice TEXT
            )
        """)

        cursor.execute("""
            CREATE TABLE IF NOT EXISTS station_registry (
                station_name TEXT PRIMARY KEY,
                crs_code TEXT
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS historical_train_runs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                timestamp TEXT,
                origin TEXT,
                destination TEXT,
                initial_delay INTEGER,
                predicted_delay INTEGER
            )
        """)
        conn.commit()

        # 2. READ AND CONVERT SYSTEM SCHEMAS LINE BY LINE Safely
        mysql_schema_path = os.path.join(KB_DIR, 'MySQL_CW2_schema.sql')
        if os.path.exists(mysql_schema_path):
            try:
                with open(mysql_schema_path, 'r', encoding='utf-8', errors='ignore') as f:
                    lines = f.readlines()
                
                cleaned_statements = []
                current_statement = ""
                
                for line in lines:
                    clean_line = line.strip()
                    # Skip database creation/management instructions completely
                    if not clean_line or clean_line.startswith('--') or clean_line.startswith('#'):
                        continue
                    if any(clean_line.upper().startswith(kw) for kw in ["CREATE DATABASE", "USE ", "DROP DATABASE"]):
                        continue
                        
                    current_statement += " " + line
                    if ';' in clean_line:
                        # Translate MySQL variables to SQLite syntax rules
                        stmt = current_statement.strip()
                        stmt = re.sub(r'(?i)\bENGINE\s*=\s*\w+', '', stmt)
                        stmt = re.sub(r'(?i)\bAUTO_INCREMENT\b', 'AUTOINCREMENT', stmt)
                        stmt = re.sub(r'(?i)\bINT\b', 'INTEGER', stmt)
                        if stmt:
                            cleaned_statements.append(stmt)
                        current_statement = ""

                for sql_stmt in cleaned_statements:
                    try:
                        cursor.execute(sql_stmt)
                    except sqlite3.OperationalError:
                        # Quietly bypass any specific non-standard SQL keywords 
                        pass
                conn.commit()
            except Exception as e:
                print(f"[Schema Core Sync Notice] Handled structural alignments safely: {e}")

        conn.close()

    def sync_and_load_stations(self):
        """Syncs structural row maps cleanly out of stations.csv into the local file db cache."""
        csv_path = os.path.join(KB_DIR, 'stations.csv')
        df = pd.DataFrame()
        if os.path.exists(csv_path):
            try:
                df = pd.read_csv(csv_path)
                conn = sqlite3.connect(DB_PATH)
                cursor = conn.cursor()
                for _, row in df.iterrows():
                    name = str(row.get('NAME', row.iloc[0])).upper().strip()
                    code = str(row.get('CRS', row.iloc[1])).upper().strip()
                    cursor.execute("""
                        INSERT OR REPLACE INTO station_registry (station_name, crs_code)
                        VALUES (?, ?)
                    """, (name, code))
                conn.commit()
                conn.close()
            except Exception:
                pass
        return df

    def load_and_mine_local_docs(self):
        """Mines local PPTM presentation decks and crawls DOCX regional management folders."""
        if not os.path.exists(KB_DIR):
            return

        pptm_files = glob.glob(os.path.join(KB_DIR, "*.pptm"))
        if Presentation:
            for filepath in pptm_files:
                try:
                    prs = Presentation(filepath)
                    filename = os.path.basename(filepath)
                    for idx, slide in enumerate(prs.slides):
                        slide_text = ""
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text += shape.text + " "
                        if slide_text.strip():
                            self.mined_knowledge_pool.append({
                                "source": f"{filename} (Slide {idx+1})",
                                "content": slide_text.strip().lower()
                            })
                except Exception:
                    pass

        if Document:
            for root, dirs, files in os.walk(KB_DIR):
                for file in files:
                    if file.endswith(".docx") and not file.startswith("~$"):
                        filepath = os.path.join(root, file)
                        try:
                            doc = Document(filepath)
                            doc_text = ""
                            for para in doc.paragraphs:
                                if para.text.strip():
                                    doc_text += para.text + "\n"
                            for table in doc.tables:
                                for row in table.rows:
                                    for cell in row.cells:
                                        if cell.text.strip():
                                            doc_text += cell.text + " "
                            if doc_text.strip():
                                self.mined_knowledge_pool.append({
                                    "source": os.path.relpath(filepath, KB_DIR),
                                    "content": doc_text.strip().lower()
                                })
                        except Exception:
                            pass

    def calculate_match_score(self, text_pool, query_tokens):
        if not text_pool or not query_tokens:
            return 0
        matches = sum(1 for token in query_tokens if token in text_pool)
        return matches / len(query_tokens)

    def evaluate_incident(self, severity, event_type, location):
        """Finds custom directive fragments from documentation file matrices using token weights."""
        search_phrase = f"{location.lower()} {event_type.lower()} {severity.lower()} disruption contingency"
        query_tokens = [w for w in re.split(r'\W+', search_phrase) if len(w) > 3]

        best_match = None
        highest_score = 0.0

        for node in self.mined_knowledge_pool:
            score = self.calculate_match_score(node["content"], query_tokens)
            if score > highest_score:
                highest_score = score
                best_match = node

        if highest_score > 0.35 and best_match:
            snippet = best_match["content"][:350].replace("\n", " ") + "..."
            return {
                "advice_staff": f"[Mined Document Plan: {best_match['source']}]\nEnact custom local protocols discovered: {snippet}",
                "advice_passenger": f"Disruption active at {location}. SWR operational squads are deploying regional backup alternatives. Please monitor platform displays."
            }

        # Backup vectors fallback
        if "FULL" in str(severity).upper():
            return {
                "advice_staff": "Enforce strict emergency mainline split. Signal immediate stop orders to all inbound active routes.",
                "advice_passenger": "Mainline rail service is suspended. Bus replacement operations are running between affected hubs."
            }
        else:
            return {
                "advice_staff": "Implement single-line bi-directional tracking configurations. Restrict speeds through affected loops to 15mph.",
                "advice_passenger": "Expect platform alterations and service extensions up to 30 minutes due to line limitations."
            }

    def log_interaction_to_db(self, mode, user_msg, bot_resp):
        try:
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO conversation_history (timestamp, session_mode, user_message, bot_response)
                VALUES (?, ?, ?, ?)
            """, (datetime.datetime.now().isoformat(), mode, user_msg, bot_resp))
            conn.commit()
            conn.close()
        except Exception:
            pass

    def log_incident_to_db(self, event_type, location, severity, staff_adv, pass_adv):
        try:
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO incident_logs (timestamp, location, event_type, severity, staff_advice, passenger_advice)
                VALUES (?, ?, ?, ?, ?, ?)
            """, (datetime.datetime.now().isoformat(), location, event_type, severity, staff_adv, pass_adv))
            conn.commit()
            conn.close()
        except Exception:
            pass

    def log_predictive_run(self, origin, destination, init_delay, pred_delay):
        try:
            conn = sqlite3.connect(DB_PATH)
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO historical_train_runs (timestamp, origin, destination, initial_delay, predicted_delay)
                VALUES (?, ?, ?, ?, ?)
            """, (datetime.datetime.now().isoformat(), origin, destination, init_delay, pred_delay))
            conn.commit()
            conn.close()
        except Exception:
            pass

expert_system = ContingencyExpertSystem()