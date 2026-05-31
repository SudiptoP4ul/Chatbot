import os
import re
import sqlite3
import glob
import json
import logging
import numpy as np
from difflib import get_close_matches
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
from groq import Groq

logging.disable(logging.CRITICAL)

BASE_DIR     = "/Users/sudiptogoldfish/code files/7059B A_AI Lab/Chatbot/knowledge based/"
DB_PATH      = os.path.join(BASE_DIR, 'swr_contingency.db')
STATIONS_CSV = os.path.join(BASE_DIR, 'stations.csv')
CHROMA_DIR   = os.path.join(BASE_DIR, 'chroma_kb')

PPTM_FILES = [
    os.path.join(BASE_DIR, '20 Batch 2 LOR 2-4 CPT (19th May 21).pptm'),
    os.path.join(BASE_DIR, 'Mainline CPT (8th May 21).pptm'),
]
DOCX_DIRS = [
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM West'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM Central'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM Metro'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM South'),
    os.path.join(BASE_DIR, 'SWR Station Disruption Plans', 'RM Waterloo'),
]

KNOWN_LABELS = [
    'Principles of Service Alteration',
    'Alternative Passenger Journey',
    'Information for Signallers',
    'Additional Information for Station Staff',
    'Passenger Information',
]
SECTION_KEYWORDS = ['off peak', 'peak', 'during peak', 'all day', 'note', 'if no access']

SWR_CONTACTS = [
    ('Bournemouth Main Line', 'Basingstoke SCP',                '50297'),
    ('Bournemouth Main Line', 'Basingstoke ASC',                '0750131'),
    ('Bournemouth Main Line', 'Basingstoke (DC) Panel',         '0750010'),
    ('Bournemouth Main Line', 'Woking SCP',                     '50341'),
    ('Bournemouth Main Line', 'Woking ASC',                     '0069410 / 0069420'),
    ('Bournemouth Main Line', 'Southampton SCP',                '075 8219'),
    ('Bournemouth Main Line', 'Southampton Ctl Supv',           '075 8107'),
    ('Bournemouth Main Line', 'Eastleigh ASC',                  '0757202'),
    ('Bournemouth Main Line', 'Bournemouth SCP',                '50375'),
    ('Bournemouth Main Line', 'Bournemouth SB',                 '0755332'),
    ('Bournemouth Main Line', 'Brockenhurst SB',                '08528311'),
    ('Bournemouth Main Line', 'Guildford SB',                   '0067562'),
    ('Bournemouth Main Line', 'Guildford SCP',                  '50359'),
    ('Bournemouth Main Line', 'Portsmouth SCP',                 '50265'),
    ('Bournemouth Main Line', 'Dorchester SB',                  '0758355'),
    ('Bournemouth Main Line', 'Salisbury SCP',                  '50412'),
    ('Bournemouth Main Line', 'GWR Train Service Control West', '08582219'),
    ('Bournemouth Main Line', 'Guildford Surbiton Supervisor',  '07771 613061'),
]

TERMINOLOGY = {
    'DM':   'Down Main - direction of travel towards London',
    'UM':   'Up Main - direction of travel away from London',
    'SL':   'Slow Line',
    'FL':   'Fast Line',
    'DMSL': 'Down Main Slow Line - slow line towards London',
    'DMFL': 'Down Main Fast Line - fast line towards London',
    'UMSL': 'Up Main Slow Line - slow line away from London',
    'UMFL': 'Up Main Fast Line - fast line away from London',
    'SWR':  'South Western Railway',
    'NR':   'Network Rail',
    'TOC':  'Train Operating Company',
    'ASC':  'Area Signalling Centre',
    'SCP':  'Station Control Point',
    'SB':   'Signal Box',
    'CSL2': 'Customer Service Level 2 - disruption response protocol',
    'CRS':  'Computer Reservation System - 3-letter station code',
}

BLOCKAGE_WORDS   = ['blockage','blocked','block','line blocked','track blocked',
                    'dmfl','dmsl','umfl','umsl','down fast','down slow','up fast','up slow',
                    'full blockage','partial blockage','one line','both lines','derailment','obstruction']
DISRUPTION_WORDS = ['disruption','disrupted','closed','closure','station problem',
                    'not stopping','trains not stopping','station disruption','station closed','platform']
PASSENGER_WORDS  = ['passenger','passengers','alternative','route','routes','bus','coach',
                    'taxi','travel','get to','how to get','alternative route',
                    'what can passengers','advice for passenger']
CONTACT_WORDS    = ['contact','phone','number','call','who do i call','telephone','ring','escalate','control']
STAFF_WORDS      = ['staff','what should i do','what do i do','instructions','principles',
                    'alteration','service alteration','divert','diversion']
SIGNALLER_WORDS  = ['signaller','signalman','signal box','information for signallers','signalling']


def contains(text, word_list):
    t = text.lower()
    return any(w in t for w in word_list)

def get_blockage_type_id(status):
    sl = status.lower()
    if 'dmfl' in sl and 'dmsl' in sl: return 5
    if 'umfl' in sl and 'umsl' in sl: return 5
    if 'dmfl' in sl: return 1
    if 'dmsl' in sl: return 2
    if 'umfl' in sl: return 3
    if 'umsl' in sl: return 4
    if 'both' in sl or ('all' in sl and 'line' in sl): return 5
    if 'one line' in sl or 'partial' in sl: return 6
    if 'station' in sl: return 7
    return 6

def detect_section(line):
    lower = line.lower().strip()
    for kw in SECTION_KEYWORDS:
        if kw in lower: return line.strip()
    return ''

def is_important(line):
    return any(w in line.lower() for w in ['must','immediately','do not','cancel','emergency','warning'])


#database builder

def build_database(db_path=DB_PATH):
    try:
        from pptx import Presentation
        from docx import Document
    except ImportError:
        os.system("pip install python-pptx python-docx -q")
        from pptx import Presentation
        from docx import Document

    if os.path.exists(db_path):
        os.remove(db_path)
    conn = sqlite3.connect(db_path)
    cur  = conn.cursor()

    cur.executescript("""
    CREATE TABLE Tiploc (
        name TEXT, longname TEXT, name_alias TEXT,
        alpha3 TEXT, tiploc TEXT PRIMARY KEY
    );
    CREATE TABLE blockage_types (id INTEGER PRIMARY KEY, description TEXT);
    CREATE TABLE ContingencyPlan (
        id INTEGER, subplan INTEGER,
        name_of_line TEXT, first_station TEXT, last_station TEXT,
        blockage_type INTEGER, blockage_status TEXT,
        PRIMARY KEY (id, subplan)
    );
    CREATE TABLE ServiceAlteration (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, alteration TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE AltPassengerJourney (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, altjourney TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE SignallerInfo (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, info TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE StationStaffInfo (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, info TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE PassengerInfo (
        cplan_id INTEGER, subplan INTEGER,
        section TEXT, info TEXT, important INTEGER DEFAULT 0
    );
    CREATE TABLE ContactDetails (
        name_of_line TEXT, contact_name TEXT, phone_numbers TEXT
    );
    CREATE TABLE StationDisruptionPlan (
        id INTEGER PRIMARY KEY, station_name TEXT,
        dateof TEXT, doc_owner TEXT, region TEXT
    );
    CREATE TABLE OtherStationSupport (
        dplan_id INTEGER, other_station_name TEXT,
        phone_numbers TEXT, time_info TEXT
    );
    CREATE TABLE AltTransportInfo (
        dplan_id INTEGER, destination TEXT,
        route_type INTEGER, route_info TEXT
    );
    CREATE TABLE StationTips (dplan_id INTEGER, issue TEXT, tip TEXT);
    """)

    cur.executemany("INSERT INTO blockage_types VALUES (?,?)", [
        (1,"DMFL - Down Main Fast Line blocked"),
        (2,"DMSL - Down Main Slow Line blocked"),
        (3,"UMFL - Up Main Fast Line blocked"),
        (4,"UMSL - Up Main Slow Line blocked"),
        (5,"Both lines blocked (full blockage)"),
        (6,"One line blocked (partial blockage)"),
        (7,"Station blocked or closed"),
    ])

    try:
        import pandas as pd
        if os.path.exists(STATIONS_CSV):
            df = pd.read_csv(STATIONS_CSV).reset_index()
            df.columns = ['name','longname','name_alias','alpha3','tiploc']
            df['name_alias'] = df['name_alias'].replace('\\N','')
            for _, row in df.iterrows():
                cur.execute("INSERT OR IGNORE INTO Tiploc VALUES (?,?,?,?,?)",
                    (row['name'],row['longname'],row['name_alias'],row['alpha3'],row['tiploc']))
    except Exception:
        pass

    cur.executemany("INSERT INTO ContactDetails VALUES (?,?,?)", SWR_CONTACTS)

    plan_count = 0
    for filepath in PPTM_FILES:
        if not os.path.exists(filepath): continue
        from pptx import Presentation
        prs = Presentation(filepath)
        for slide in prs.slides:
            title = ""; table_cells = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape,"text") and shape.text.strip():
                        if not title: title = shape.text.strip()
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                t = cell.text.strip()
                                if t: table_cells.append(t)
                except Exception: pass
            if not re.match(r'Plan\s+[\d\.a]+', title): continue
            title_clean = title.replace('\t',' ').replace('\x0b',' ').strip()
            m = re.match(r'Plan\s+([\d]+)[\.a]?([\d]*)\s+(.*)', title_clean)
            if not m: continue
            plan_id = int(m.group(1)); subplan = int(m.group(2)) if m.group(2) else 0
            loc_str = m.group(3).strip(); parts = re.split(r'\s+-\s+', loc_str)
            first_st = parts[0].strip() if parts else ""
            last_st  = parts[1].strip() if len(parts) > 1 else ""
            line = ""; status = ""; field_map = {}
            for idx, cell in enumerate(table_cells):
                if ('Main Line' in cell or 'Line' in cell) and len(cell) < 60:
                    line = cell.replace('\n',' ').strip()
                elif cell.startswith('Status:'):
                    status = cell.replace('Status:','').strip()
                elif cell in KNOWN_LABELS and idx+1 < len(table_cells):
                    field_map[cell] = table_cells[idx+1]
            cur.execute("INSERT OR IGNORE INTO ContingencyPlan VALUES (?,?,?,?,?,?,?)",
                (plan_id,subplan,line,first_st,last_st,get_blockage_type_id(status),status))

            def insert_bullets(table, col, text, pid, sp):
                if not text or text == 'N/A': return
                current_section = ''
                for ln in text.split('\n'):
                    ln = ln.strip().lstrip('*').lstrip('-').lstrip('o').strip()
                    if not ln: continue
                    sec = detect_section(ln)
                    if sec: current_section = sec; continue
                    cur.execute(f"INSERT INTO {table} VALUES (?,?,?,?,?)",
                                (pid,sp,current_section,ln,1 if is_important(ln) else 0))

            insert_bullets("ServiceAlteration","alteration",
                field_map.get('Principles of Service Alteration',''),plan_id,subplan)
            insert_bullets("AltPassengerJourney","altjourney",
                field_map.get('Alternative Passenger Journey',''),plan_id,subplan)
            insert_bullets("SignallerInfo","info",
                field_map.get('Information for Signallers',''),plan_id,subplan)
            insert_bullets("StationStaffInfo","info",
                field_map.get('Additional Information for Station Staff',''),plan_id,subplan)
            insert_bullets("PassengerInfo","info",
                field_map.get('Passenger Information',''),plan_id,subplan)
            plan_count += 1

    docx_files = []
    for d in DOCX_DIRS:
        if os.path.exists(d):
            docx_files += glob.glob(os.path.join(d,"*.docx"))
    docx_files += glob.glob(os.path.join(BASE_DIR,"Station_Disruption_Plan_*.docx"))
    docx_files = sorted(set(docx_files))

    dplan_id = 1
    from docx import Document
    for filepath in docx_files:
        try:
            doc = Document(filepath)
            station_name = ""; doc_owner = ""; dateof = ""; region = "Unknown"
            for r in ["RM West","RM Central","RM Metro","RM South","RM Waterloo"]:
                if r in filepath: region = r; break
            if region == "Unknown": region = "RM West"
            for para in doc.paragraphs:
                text = para.text.strip()
                if text and para.style.name.startswith("Heading") and not station_name:
                    station_name = text.replace("Station","").strip(); break
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if len(cells) >= 2:
                        if "Document Owner" in cells[0]: doc_owner = cells[1]
                        if "Date" in cells[0] and not dateof: dateof = cells[1]
            if not station_name:
                fn = os.path.basename(filepath)
                m2 = re.search(r'Plan_-_(.+?)_Issue', fn)
                if m2: station_name = m2.group(1).replace('_',' ').strip()
            cur.execute("INSERT INTO StationDisruptionPlan VALUES (?,?,?,?,?)",
                        (dplan_id,station_name,dateof,doc_owner,region))
            for table in doc.tables:
                headers = [c.text.strip().upper() for c in table.rows[0].cells]
                if any("ROUTE AVAILABLE" in h for h in headers):
                    for row in table.rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        dest = cells[1] if len(cells) > 1 else ""
                        if not dest or dest in ("Destination",""): continue
                        rail    = cells[0] if len(cells) > 0 else ""
                        bus     = cells[2] if len(cells) > 2 else ""
                        walking = cells[3] if len(cells) > 3 else ""
                        if rail: cur.execute("INSERT INTO AltTransportInfo VALUES (?,?,?,?)",(dplan_id,dest,0,rail))
                        if bus:  cur.execute("INSERT INTO AltTransportInfo VALUES (?,?,?,?)",(dplan_id,dest,1,bus))
                        if walking and walking.lower() not in ('n/a',''):
                            cur.execute("INSERT INTO AltTransportInfo VALUES (?,?,?,?)",(dplan_id,dest,2,walking))
                if len(headers) >= 2 and "ISSUE" in headers[0]:
                    for row in table.rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        if cells[0] and cells[0].upper() not in ("ISSUE",""):
                            cur.execute("INSERT INTO StationTips VALUES (?,?,?)",
                                (dplan_id,cells[0][:255],cells[1][:500] if len(cells)>1 else ""))
                if any("STATION NAME" in h or "TELEPHONE" in h for h in headers):
                    for row in table.rows[1:]:
                        cells = [c.text.strip() for c in row.cells]
                        if cells[0] and cells[0] not in ("N/A","Station Name",""):
                            cur.execute("INSERT INTO OtherStationSupport VALUES (?,?,?,?)",
                                (dplan_id,cells[0],cells[1] if len(cells)>1 else "",cells[2] if len(cells)>2 else ""))
            dplan_id += 1
        except Exception: pass

    conn.commit(); conn.close()


#database access 

class ContingencyDB:
    def __init__(self, db_path=DB_PATH):
        if not os.path.exists(db_path):
            build_database(db_path)
        self.conn = sqlite3.connect(db_path)
        self.conn.row_factory = sqlite3.Row
        self._build_search_index()

    def _build_search_index(self):
        cur = self.conn.cursor()
        cur.execute("SELECT id,subplan,first_station,last_station,blockage_status FROM ContingencyPlan")
        rows = cur.fetchall()
        self.cp_ids    = [(r['id'],r['subplan']) for r in rows]
        self.cp_corpus = [f"{r['first_station']} {r['last_station']} {r['blockage_status']}" for r in rows]
        for i,(pid,sp) in enumerate(self.cp_ids):
            cur.execute("SELECT alteration FROM ServiceAlteration WHERE cplan_id=? AND subplan=?",(pid,sp))
            self.cp_corpus[i] += " " + " ".join(r['alteration'] for r in cur.fetchall())
        self.cp_vectorizer = TfidfVectorizer(stop_words='english',ngram_range=(1,2),sublinear_tf=True)
        self.cp_matrix     = self.cp_vectorizer.fit_transform(self.cp_corpus)
        cur.execute("SELECT id,station_name FROM StationDisruptionPlan")
        rows = cur.fetchall()
        self.dp_ids   = [r['id']           for r in rows]
        self.dp_names = [r['station_name'] for r in rows]
        self.cp_stations = sorted(set(
            s[0] for s in self.conn.execute("SELECT first_station FROM ContingencyPlan") if s[0]
        ) | set(
            s[0] for s in self.conn.execute("SELECT last_station FROM ContingencyPlan") if s[0]
        ))
        self.dp_stations = sorted(set(s for s in self.dp_names if s))
        self.tiploc_cache = {}
        for row in self.conn.execute("SELECT name,alpha3,tiploc FROM Tiploc"):
            self.tiploc_cache[row['name'].upper()] = {'crs':row['alpha3'],'tiploc':row['tiploc']}

    def fuzzy_match(self, raw, station_list):
        if not raw: return None
        raw = raw.strip().title()
        if raw in station_list: return raw
        close = get_close_matches(raw, station_list, n=1, cutoff=0.5)
        if close: return close[0]
        for s in station_list:
            if raw.lower() in s.lower(): return s
        return None

    def search_contingency(self, query, loc_from='', loc_to='', top_n=3):
        full_query = f"{loc_from} {loc_to} {query}".strip()
        q_vec  = self.cp_vectorizer.transform([full_query])
        scores = cosine_similarity(q_vec, self.cp_matrix)[0]
        loc_from_m = self.fuzzy_match(loc_from, self.cp_stations) if loc_from else None
        loc_to_m   = self.fuzzy_match(loc_to,   self.cp_stations) if loc_to   else None
        cur = self.conn.cursor(); boosted = []
        for i,(pid,sp) in enumerate(self.cp_ids):
            score = float(scores[i])
            cur.execute("SELECT first_station,last_station FROM ContingencyPlan WHERE id=? AND subplan=?",(pid,sp))
            row = cur.fetchone()
            if row:
                if loc_from_m and row['first_station'].lower() == loc_from_m.lower(): score += 0.4
                if loc_to_m   and row['last_station'].lower()  == loc_to_m.lower():   score += 0.4
            boosted.append((pid,sp,score))
        boosted.sort(key=lambda x: x[2], reverse=True)
        results = []
        for pid,sp,score in boosted[:top_n]:
            if score > 0.01:
                plan = self.get_full_plan(pid,sp,score)
                if plan: results.append(plan)
        return results

    def search_station_disruption(self, station_name):
        matched = self.fuzzy_match(station_name, self.dp_stations)
        if matched:
            cur = self.conn.cursor()
            cur.execute("SELECT id FROM StationDisruptionPlan WHERE station_name=?",(matched,))
            row = cur.fetchone()
            if row: return self.get_disruption_plan(row['id'])
        return None

    def get_full_plan(self, plan_id, subplan, score=0):
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM ContingencyPlan WHERE id=? AND subplan=?",(plan_id,subplan))
        cp = cur.fetchone()
        if not cp: return None
        def get_bullets(table, col):
            cur.execute(f"SELECT section,{col},important FROM {table} "
                        f"WHERE cplan_id=? AND subplan=? ORDER BY rowid",(plan_id,subplan))
            return [{'text':r[1],'section':r[0],'important':bool(r[2])} for r in cur.fetchall()]
        return {
            "plan_id":plan_id,"subplan":subplan,"score":score,"line":cp['name_of_line'],
            "first_station":cp['first_station'],"last_station":cp['last_station'],
            "blockage_status":cp['blockage_status'],
            "principles":  get_bullets("ServiceAlteration",  "alteration"),
            "alt_journey": get_bullets("AltPassengerJourney","altjourney"),
            "signaller":   get_bullets("SignallerInfo",       "info"),
            "staff":       get_bullets("StationStaffInfo",   "info"),
            "passenger":   get_bullets("PassengerInfo",      "info"),
        }

    def get_disruption_plan(self, dplan_id):
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM StationDisruptionPlan WHERE id=?",(dplan_id,))
        dp = cur.fetchone()
        if not dp: return None
        cur.execute("SELECT destination,route_type,route_info FROM AltTransportInfo "
                    "WHERE dplan_id=? ORDER BY destination,route_type",(dplan_id,))
        transport = {}
        for row in cur.fetchall():
            dest = row['destination']
            if dest not in transport:
                transport[dest] = {'destination':dest,'rail':'','bus':'','walking':''}
            if   row['route_type']==0: transport[dest]['rail']    = row['route_info']
            elif row['route_type']==1: transport[dest]['bus']     = row['route_info']
            elif row['route_type']==2: transport[dest]['walking'] = row['route_info']
        cur.execute("SELECT * FROM StationTips         WHERE dplan_id=?",(dplan_id,))
        tips    = [dict(r) for r in cur.fetchall()]
        cur.execute("SELECT * FROM OtherStationSupport WHERE dplan_id=?",(dplan_id,))
        support = [dict(r) for r in cur.fetchall()]
        return {"station_name":dp['station_name'],"dateof":dp['dateof'],
                "region":dp['region'],"doc_owner":dp['doc_owner'],
                "transport":list(transport.values()),"tips":tips,"support":support}

    def get_contacts(self, line_name=None):
        cur = self.conn.cursor()
        if line_name:
            cur.execute("SELECT contact_name,phone_numbers FROM ContactDetails "
                        "WHERE name_of_line LIKE ?",(f'%{line_name}%',))
        else:
            cur.execute("SELECT contact_name,phone_numbers FROM ContactDetails")
        return [dict(r) for r in cur.fetchall()]

    def lookup_tiploc(self, station_name):
        return self.tiploc_cache.get(station_name.upper().strip(), None)

    def get_all_cp_stations(self): return self.cp_stations
    def get_all_dp_stations(self): return self.dp_stations
    def explain_term(self, term):  return TERMINOLOGY.get(term.upper().strip(), None)


#RAG knowledge base 

class RAGKnowledgeBase:
    def __init__(self, db):
        self.db=db; self.client=None; self.coll=None; self._ready=False
        try:
            import chromadb
            from chromadb.utils import embedding_functions
            self.client = chromadb.PersistentClient(path=CHROMA_DIR)
            ef = embedding_functions.DefaultEmbeddingFunction()
            self.coll = self.client.get_or_create_collection(
                name="swr_contingency_plans",embedding_function=ef,
                metadata={"hnsw:space":"cosine"})
            if self.coll.count() == 0: self._index_plans()
            self._ready = True
        except ImportError: pass
        except Exception:   pass

    def _index_plans(self):
        conn=self.db.conn; cur=conn.cursor()
        cur.execute("SELECT id,subplan,name_of_line,first_station,last_station,blockage_status FROM ContingencyPlan")
        plans=cur.fetchall(); docs=[]; ids=[]; metas=[]
        section_map={
            "ServiceAlteration":  ("alteration","Principles of Service Alteration"),
            "AltPassengerJourney":("altjourney", "Alternative Passenger Journey"),
            "SignallerInfo":      ("info",        "Information for Signallers"),
            "StationStaffInfo":   ("info",        "Additional Information for Station Staff"),
            "PassengerInfo":      ("info",        "Passenger Information"),
        }
        for plan in plans:
            pid,sp=plan['id'],plan['subplan']
            header=(f"Plan {pid}.{sp} | {plan['name_of_line']} | "
                    f"{plan['first_station']} to {plan['last_station']} | "
                    f"Status: {plan['blockage_status']}")
            for table,(col,label) in section_map.items():
                cur.execute(f"SELECT {col} FROM {table} WHERE cplan_id=? AND subplan=? ORDER BY rowid",(pid,sp))
                rows=cur.fetchall()
                if not rows: continue
                body="\n".join(r[0] for r in rows if r[0])
                if not body.strip(): continue
                chunk=f"{header}\n\n[{label}]\n{body}"; uid=f"plan_{pid}_{sp}_{table}"
                docs.append(chunk); ids.append(uid)
                metas.append({"plan_id":pid,"subplan":sp,"section":label,
                               "first_station":plan['first_station'],
                               "last_station":plan['last_station'],"status":plan['blockage_status']})
        cur.execute("SELECT id,station_name,region FROM StationDisruptionPlan")
        for dp in cur.fetchall():
            did=dp['id']
            cur.execute("SELECT destination,route_info FROM AltTransportInfo WHERE dplan_id=?",(did,))
            transport="\n".join(f"To {r['destination']}: {r['route_info']}" for r in cur.fetchall())
            cur.execute("SELECT issue,tip FROM StationTips WHERE dplan_id=?",(did,))
            tips="\n".join(f"Issue: {r['issue']} - Tip: {r['tip']}" for r in cur.fetchall())
            body=f"Station: {dp['station_name']} ({dp['region']})\n"
            if transport: body+=f"\nAlternative Transport:\n{transport}"
            if tips:      body+=f"\nStation Tips:\n{tips}"
            docs.append(body); ids.append(f"dplan_{did}")
            metas.append({"plan_id":did,"subplan":0,"section":"Station Disruption",
                           "first_station":dp['station_name'],"last_station":"","status":"station disruption"})
        self._text_store={}
        for uid,doc in zip(ids,docs): self._text_store[uid]=doc
        for i in range(0,len(docs),500):
            self.coll.add(documents=ids[i:i+500],ids=ids[i:i+500],metadatas=metas[i:i+500])
        store_path=os.path.join(CHROMA_DIR,'text_store.json')
        os.makedirs(CHROMA_DIR,exist_ok=True)
        with open(store_path,'w') as f: json.dump(self._text_store,f)

    def _load_text_store(self):
        store_path=os.path.join(CHROMA_DIR,'text_store.json')
        self._text_store=json.load(open(store_path)) if os.path.exists(store_path) else {}

    def retrieve(self, query, n_results=5):
        if not self._ready or self.coll is None: return []
        if not hasattr(self,'_text_store'): self._load_text_store()
        try:
            results=self.coll.query(query_texts=[query],n_results=min(n_results,self.coll.count()))
            return [self._text_store.get(uid,'') for uid in (results['ids'][0] if results['ids'] else []) if uid]
        except Exception: return []

    @property
    def ready(self): return self._ready


#singletons

_db = None; _rag_kb = None

def get_db():
    global _db
    if _db is None: _db = ContingencyDB(DB_PATH)
    return _db

def get_rag_kb():
    global _rag_kb
    if _rag_kb is None: _rag_kb = RAGKnowledgeBase(get_db())
    return _rag_kb


#Groq wrapper


GROQ_MODEL = "llama-3.1-8b-instant"   # 500k TPD — good balance of quality and cost


def call_groq(system_prompt, messages, max_tokens=600):
    api_key = os.environ.get("GROQ_API_KEY","")
    if not api_key:
        print("[Groq] ERROR: GROQ_API_KEY is not set")
        return ""
    try:
        client = Groq(api_key=api_key)
        response = client.chat.completions.create(
            model=GROQ_MODEL,
            max_tokens=max_tokens,
            messages=[{"role":"system","content":system_prompt}] + messages,
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        print(f"[Groq] ERROR: {type(e).__name__}: {e}")
        return ""


#data formatters

def format_contacts(db):
    lines = []
    for c in db.get_contacts():
        lines.append(f"  {c['contact_name']:<40} {c['phone_numbers']}")
    return "\n".join(lines)

def format_bullets(bullets):
    lines = []
    current_section = ''
    for b in bullets:
        if b['section'] and b['section'] != current_section:
            current_section = b['section']
            lines.append(f"  [{current_section}]")
        prefix = "  !  " if b['important'] else "  -  "
        lines.append(f"{prefix}{b['text']}")
    return "\n".join(lines)

def format_contingency_plan(plan, section='all'):
    lines = []
    lines.append(f"Plan {plan['plan_id']}.{plan['subplan']}  |  {plan['first_station']} to {plan['last_station']}")
    lines.append(f"Status : {plan['blockage_status']}")
    if plan.get('line'): lines.append(f"Line   : {plan['line']}")
    lines.append("")
    if section in ('all','staff') and plan['principles']:
        lines.append("Principles of Service Alteration:")
        lines.append(format_bullets(plan['principles']))
        lines.append("")
    if section in ('all','staff') and plan['staff']:
        lines.append("Additional Instructions for Station Staff:")
        lines.append(format_bullets(plan['staff']))
        lines.append("")
    if section in ('all','signaller') and plan['signaller']:
        lines.append("Information for Signallers:")
        lines.append(format_bullets(plan['signaller']))
        lines.append("")
    elif section == 'signaller' and not plan['signaller']:
        lines.append("No signaller-specific instructions found in this plan.")
        lines.append("")
    if section in ('all','passenger') and plan['alt_journey']:
        lines.append("Alternative Passenger Journey:")
        lines.append(format_bullets(plan['alt_journey']))
        lines.append("")
    if section in ('all','passenger') and plan['passenger']:
        lines.append("Passenger Information:")
        lines.append(format_bullets(plan['passenger']))
        lines.append("")
    return "\n".join(lines)

def format_disruption_plan(plan, section='all'):
    lines = []
    lines.append(f"Station Disruption Plan: {plan['station_name']}")
    if plan['region']:   lines.append(f"Region : {plan['region']}")
    if plan['dateof']:   lines.append(f"Date   : {plan['dateof']}")
    lines.append("")
    if section in ('all','passenger') and plan['transport']:
        lines.append("Alternative Travel Options:")
        for t in plan['transport']:
            if t['destination']:
                lines.append(f"  To {t['destination']}:")
                if t['rail']:    lines.append(f"    Rail    - {t['rail']}")
                if t['bus']:     lines.append(f"    Bus     - {t['bus']}")
                if t['walking']: lines.append(f"    Walking - {t['walking']}")
        lines.append("")
    if section in ('all','staff') and plan['tips']:
        lines.append("Station Tips:")
        for t in plan['tips']:
            if t['issue']:
                lines.append(f"  Issue : {t['issue'][:120]}")
                if t['tip']: lines.append(f"  Tip   : {t['tip'][:120]}")
        lines.append("")
    if section in ('all','contacts') and plan['support']:
        lines.append("Nearby Station Contacts:")
        for s in plan['support']:
            if s['other_station_name']:
                line = f"  {s['other_station_name']}"
                if s['phone_numbers']: line += f"   Tel: {s['phone_numbers']}"
                lines.append(line)
        lines.append("")
    return "\n".join(lines)


# DB context builder

def pull_db_context(db, rag_kb, user_text, history):
    """
    """
    context_parts = []

    # build a combined view of what has been said so far
    full_history_text = " ".join(m['content'] for m in history).lower() + " " + user_text.lower()

    # extract locations from the current message
    # "between X and Y" / "from X to Y"
    pair = re.search(
        r'(?:between|from)\s+([A-Z][a-zA-Z &\-]+?)\s+(?:and|to)\s+([A-Z][a-zA-Z &\-]+?)(?=\s|,|$)',
        user_text
    )
    at_match = re.search(
        r'\bat\s+([A-Z][a-zA-Z &\-]+?)(?:\s+[Ss]tation)?\s*(?:$|,|\.|\bis\b|\bhas\b)',
        user_text
    )
    for_match = re.search(
        r'\bfor\s+([A-Z][a-zA-Z &\-]+?)(?:\s+[Ss]tation)?\s*(?:$|,|\.)',
        user_text
    )
    # line codes
    code_match = re.search(r'\b(DMFL|DMSL|UMFL|UMSL)\b', user_text, re.IGNORECASE)

    station_from = (pair.group(1).strip() if pair
                    else at_match.group(1).strip() if at_match
                    else for_match.group(1).strip() if for_match
                    else None)
    station_to   = pair.group(2).strip() if pair else None
    blockage_code = code_match.group(1).upper() if code_match else None

    # also try to find stations mentioned anywhere in the current message
    if not station_from:
        all_stations = db.get_all_cp_stations() + db.get_all_dp_stations()
        for s in all_stations:
            if s.lower() in user_text.lower():
                if not station_from:
                    station_from = s
                elif not station_to and s != station_from:
                    station_to = s

    # terminology lookups 
    for abbrev, meaning in TERMINOLOGY.items():
        if re.search(rf'\b{abbrev}\b', user_text, re.IGNORECASE):
            context_parts.append(f"TERMINOLOGY: {abbrev} = {meaning}")

    # contact numbers
    if contains(full_history_text, CONTACT_WORDS):
        ctx = format_contacts(db)
        if ctx:
            context_parts.append("SWR CONTROL CONTACTS:\n" + ctx)

    # work out which section the user needs 
    if contains(full_history_text, SIGNALLER_WORDS):
        section = 'signaller'
    elif contains(full_history_text, PASSENGER_WORDS) and not contains(full_history_text, SIGNALLER_WORDS):
        section = 'passenger'
    elif contains(full_history_text, STAFF_WORDS) and not contains(full_history_text, SIGNALLER_WORDS):
        section = 'staff'
    else:
        section = 'all'

    # contingency plan lookup 
    is_blockage = contains(full_history_text, BLOCKAGE_WORDS) or blockage_code is not None
    if station_from and is_blockage:
        query = f"{blockage_code or 'blockage'} {station_from} {station_to or ''}"
        results = db.search_contingency(query, station_from, station_to or '', top_n=2)
        for plan in results:
            context_parts.append(
                f"CONTINGENCY PLAN DATA:\n{format_contingency_plan(plan, section)}"
            )

    # station disruption plan lookup
    is_disruption = contains(full_history_text, DISRUPTION_WORDS)
    if station_from and is_disruption:
        plan = db.search_station_disruption(station_from)
        if plan:
            context_parts.append(
                f"STATION DISRUPTION PLAN DATA:\n{format_disruption_plan(plan, section)}"
            )

    # RAG semantic search
    # runs as a top-up when structured lookups found something, or as the
    # primary source when they didn't (e.g. vague or conversational messages)
    if rag_kb.ready and (station_from or is_blockage or is_disruption):
        chunks = rag_kb.retrieve(user_text, n_results=3)
        if chunks:
            context_parts.append(
                "ADDITIONAL KNOWLEDGE BASE CONTEXT:\n" + "\n\n---\n\n".join(chunks)
            )

    # ── station codes as a reference footer
    if station_from:
        tiploc = db.lookup_tiploc(station_from.upper())
        if tiploc:
            context_parts.append(
                f"STATION CODE: {station_from}   CRS: {tiploc['crs']}   TIPLOC: {tiploc['tiploc']}"
            )
    if station_to:
        tiploc2 = db.lookup_tiploc(station_to.upper())
        if tiploc2:
            context_parts.append(
                f"STATION CODE: {station_to}   CRS: {tiploc2['crs']}   TIPLOC: {tiploc2['tiploc']}"
            )

    return "\n\n".join(context_parts)


#system prompt

def build_system_prompt(user_type, db_context):
    """
    Builds the system prompt fresh each turn with the current DB context injected.

    The whole approach here is: give the LLM a clearly defined role, the right
    data, and good instincts about tone — then get out of the way. The LLM
    handles all the natural language variation. We just keep it grounded.
    """
    if user_type == 'passenger':
        role = (
            "You are a helpful travel assistant for South Western Railway passengers.\n"
            "Use plain, simple English. No jargon, no plan numbers, no operational detail.\n"
            "Focus on getting the passenger where they need to go. Be warm and reassuring.\n"
            "If you do not have enough information in the context, apologise briefly and "
            "direct them to National Rail Enquiries (nationalrail.co.uk) or station staff."
        )
    else:
        role = (
            "You are the SWR Contingency Assistant, a support system for South Western Railway "
            "operational staff during live incidents.\n\n"
            "You help with: line blockages, station disruptions, signaller instructions, "
            "passenger advice, staff procedures, contact numbers, and railway terminology.\n\n"
            "When context contains plan data, use it directly and include the plan number.\n"
            "When the user asks for a specific section (e.g. signaller info, passenger advice), "
            "give only that section — do not dump the entire plan.\n"
            "If information is missing from the context, say so and suggest contacting Network Control.\n"
            "Keep responses concise — staff are dealing with live incidents."
        )

    style = (
        "\nTone and style:\n"
        "- Talk like a knowledgeable colleague. Not a system, not a helpdesk script.\n"
        "- Do not open with 'Certainly!', 'Of course!', 'Great question!' or anything like that.\n"
        "- Do not say 'As an AI...' or describe your own capabilities unprompted.\n"
        "- If someone asks something unrelated to railways, respond briefly and warmly, "
        "then bring things back to how you can help.\n"
        "- Do not end every message with 'Is there anything else I can help with?' — "
        "only ask a follow-up if you genuinely need more information to help.\n"
        "- If you need to ask something, ask one thing only. Never stack questions.\n"
        "- Keep responses focused. The person is busy."
    )

    if db_context:
        data_block = (
            "\n\nThe following data has been retrieved from the SWR contingency plan database. "
            "Use this as your primary source. Do not invent details beyond what is here.\n\n"
            + db_context
        )
    else:
        data_block = (
            "\n\nNo specific plan data has been retrieved yet. "
            "If the person is describing an incident, ask which station or line is affected "
            "so you can pull up the right plan."
        )

    return (role + style + data_block).strip()


#main conversation loop

def chat(user_type='staff'):
    """
    One continuous loop. No branching state machine, no rigid handler chains.

    Each turn:
      1. User types something.
      2. We pull whatever DB data is relevant to that message.
      3. We build a fresh system prompt with that data injected as context.
      4. We pass the full conversation history to Groq.
      5. The LLM responds naturally, grounded in the plan data.
      6. Repeat.

    The LLM handles all the language variation — greetings, follow-ups,
    clarifications, off-topic messages, terminology questions, everything.
    We just make sure it always has the right data in front of it.
    """
    db     = get_db()
    rag_kb = get_rag_kb()
    history = []

    print()
    if user_type == 'passenger':
        print("Assistant: Hi, what's going on — where are you trying to get to?")
    else:
        print("Assistant: Ready. What's the incident?")
    print()

    while True:
        try:
            user_input = input("You: ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\nAssistant: Take care.")
            break

        if not user_input:
            continue

        if user_input.lower() in ('exit','quit','bye','q'):
            print("\nAssistant: Take care.")
            break

        # pull relevant DB context for this specific message
        db_context = pull_db_context(db, rag_kb, user_input, history)

        # build the system prompt with that context baked in
        system = build_system_prompt(user_type, db_context)

        # add the user message to history before calling
        history.append({"role":"user","content":user_input})

        # call the LLM with the full conversation so far
        response = call_groq(system, history, max_tokens=600)

        if not response:
            response = ("I'm having trouble reaching the system right now. "
                        "Please contact Network Control directly.")

        print(f"\nAssistant: {response}\n")

        # store the response so context carries forward into the next turn
        history.append({"role":"assistant","content":response})


#startup 

def ask_user_type():
    print()
    print("  Are you a member of staff or a passenger?")
    print()

    staff_words     = ['staff','operator','signaller','driver','controller','control','manager',
                       'employee','worker','dispatcher','train crew','crew','network','rail worker','1','s']
    passenger_words = ['passenger','traveller','traveler','customer','commuter','public','user',
                       'visitor','tourist','i am travelling','i need to get','2','p']

    while True:
        raw = input("You: ").strip().lower()
        if not raw: continue
        if any(w in raw for w in staff_words):     return 'staff'
        if any(w in raw for w in passenger_words): return 'passenger'

        result = call_groq(
            "Classify the speaker as 'staff' (railway employee) or 'passenger' (traveller). "
            "Return only one word: staff or passenger.",
            [{"role":"user","content":raw}], max_tokens=10
        )
        r = result.strip().lower()
        if 'staff'     in r: return 'staff'
        if 'passenger' in r: return 'passenger'

        print("\n  Just say 'staff' or 'passenger' and we'll get started.\n")


def main():
    if not os.environ.get("GROQ_API_KEY",""):
        print("\nError: GROQ_API_KEY is not set.\n")
        return
    get_db()
    get_rag_kb()
    user_type = ask_user_type()
    chat(user_type)

# GUI integration 

class ContingencySession:
    """
    Holds the conversation history and user_type for one chat session.
    Create one instance per GUI session (or per user if multi-user).
    Call .process(user_input) each turn; it returns the assistant reply string.
    """
    def __init__(self, user_type: str = 'staff'):
        self.user_type = user_type
        self.history: list = []
        self._db     = None
        self._rag_kb = None

    def _get_db(self):
        if self._db is None:
            self._db = get_db()
        return self._db

    def _get_rag(self):
        if self._rag_kb is None:
            self._rag_kb = get_rag_kb()
        return self._rag_kb

    def process(self, user_input: str) -> str:
        """Send one message, get one reply. History is maintained internally."""
        try:
            db     = self._get_db()
            rag_kb = self._get_rag()

            db_context = pull_db_context(db, rag_kb, user_input, self.history)
            system     = build_system_prompt(self.user_type, db_context)

            self.history.append({"role": "user", "content": user_input})
            response = call_groq(system, self.history, max_tokens=600)

            if not response:
                response = ("I'm having trouble reaching the system right now. "
                            "Please contact Network Control directly.")

            self.history.append({"role": "assistant", "content": response})
            return response

        except Exception as e:
            import traceback
            traceback.print_exc()
            return (f"Something went wrong: {type(e).__name__}: {e}. "
                    f"Please contact Network Control directly.")

    def reset(self):
        """Clear history to start a fresh conversation."""
        self.history = []


_sessions: dict = {}

def process_contingency_query(user_input: str, user_type: str = 'staff') -> str:
    """
    Stateful per-user-type wrapper.
    For GUI use, prefer creating a ContingencySession instance directly
    so each chat session has its own isolated history.
    """
    global _sessions
    if user_type not in _sessions:
        _sessions[user_type] = ContingencySession(user_type=user_type)
    return _sessions[user_type].process(user_input)


def reset_contingency_session(user_type: str = 'staff'):
    """Call this when the GUI user ends their session or switches mode."""
    global _sessions
    if user_type in _sessions:
        _sessions[user_type].reset()


def process_contingency_query_safe(user_input: str, user_type: str = 'staff') -> str:
    """Safe wrapper — never raises, always returns a string."""
    try:
        return process_contingency_query(user_input, user_type)
    except Exception as e:
        print(f"[Expert System Error] {e}")
        return f"Something went wrong. Please contact Network Control directly."

if __name__ == "__main__":
    main()
